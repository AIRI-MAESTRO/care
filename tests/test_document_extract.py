"""Unit tests for ``care.runtime.document_extract`` — plain-text extraction
from office / rich-text documents behind the chat surface's ``@``-file refs.

Each format test builds a tiny *real* document with its backing library
(``importorskip`` so a stripped environment skips rather than errors) and
asserts the extracted text round-trips. The dispatch + error paths need no
third-party libraries.
"""

from __future__ import annotations

import builtins

import pytest

from care.runtime.document_extract import (
    DOCUMENT_EXTENSIONS,
    DocumentExtractionError,
    extract_document_text,
)


class TestDispatch:
    def test_supported_extensions_cover_office_formats(self):
        assert {".docx", ".pptx", ".xlsx", ".odt", ".rtf"} <= DOCUMENT_EXTENSIONS

    def test_unsupported_extension_raises(self, tmp_path):
        path = tmp_path / "note.xyz"
        path.write_text("hi")
        with pytest.raises(DocumentExtractionError, match="unsupported"):
            extract_document_text(path)

    def test_legacy_doc_hints_modern_format(self, tmp_path):
        # Extension is checked before any read, so the bytes don't matter.
        path = tmp_path / "old.doc"
        path.write_bytes(b"\xd0\xcf\x11\xe0")
        with pytest.raises(DocumentExtractionError, match=r"re-save as `\.docx`"):
            extract_document_text(path)

    def test_legacy_xls_hints_xlsx(self, tmp_path):
        path = tmp_path / "old.xls"
        path.write_bytes(b"\xd0\xcf\x11\xe0")
        with pytest.raises(DocumentExtractionError, match=r"re-save as `\.xlsx`"):
            extract_document_text(path)


class TestDocx:
    def test_paragraphs_and_tables(self, tmp_path):
        docx = pytest.importorskip("docx")
        path = tmp_path / "sample.docx"
        document = docx.Document()
        document.add_paragraph("Hello world")
        document.add_paragraph("Second line")
        table = document.add_table(rows=1, cols=2)
        table.rows[0].cells[0].text = "Alpha"
        table.rows[0].cells[1].text = "Beta"
        document.save(str(path))

        text = extract_document_text(path)
        assert "Hello world" in text
        assert "Second line" in text
        assert "Alpha | Beta" in text


class TestPptx:
    def test_slides_extracted_with_markers(self, tmp_path):
        pptx = pytest.importorskip("pptx")
        from pptx.util import Inches

        path = tmp_path / "deck.pptx"
        presentation = pptx.Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # blank
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = "Quarterly review"
        presentation.save(str(path))

        text = extract_document_text(path)
        assert "Quarterly review" in text
        assert "--- slide 1 ---" in text


class TestXlsx:
    def test_cells_extracted_per_sheet(self, tmp_path):
        openpyxl = pytest.importorskip("openpyxl")
        path = tmp_path / "book.xlsx"
        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.title = "Data"
        sheet.append(["name", "score"])
        sheet.append(["alpha", 42])
        workbook.save(str(path))

        text = extract_document_text(path)
        assert "--- sheet Data ---" in text
        assert "name | score" in text
        assert "alpha | 42" in text


class TestOdt:
    def test_paragraphs_extracted(self, tmp_path):
        pytest.importorskip("odf")
        from odf.opendocument import OpenDocumentText
        from odf.text import P

        path = tmp_path / "doc.odt"
        document = OpenDocumentText()
        document.text.addElement(P(text="OpenDocument body text"))
        document.save(str(path))

        text = extract_document_text(path)
        assert "OpenDocument body text" in text


class TestRtf:
    def test_strips_control_words(self, tmp_path):
        pytest.importorskip("striprtf")
        path = tmp_path / "note.rtf"
        path.write_text(
            r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Times;}}\f0 Plain RTF text.\par}"
        )
        text = extract_document_text(path)
        assert "Plain RTF text." in text


class TestMissingLibrary:
    def test_missing_lib_surfaces_install_hint(self, tmp_path, monkeypatch):
        """A stripped install (backing lib absent) raises a friendly
        ``pip install`` hint, not a bare ImportError."""
        real_import = builtins.__import__

        def fake_import(name, *args, **kwargs):
            if name == "docx":
                raise ImportError("simulated missing python-docx")
            return real_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", fake_import)
        path = tmp_path / "x.docx"
        path.write_bytes(b"PK\x03\x04")  # never read — the import fails first
        with pytest.raises(
            DocumentExtractionError, match=r"pip install python-docx"
        ):
            extract_document_text(path)
