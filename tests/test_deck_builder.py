"""Deterministic ``.pptx`` builder (:mod:`care.runtime.deck_builder`).

The model-independent fallback: turn generated slide text into a real openable
deck. No LLM / Docker / network.
"""

from __future__ import annotations

from pathlib import Path

import pytest

from care.runtime.deck_builder import (
    _split_into_slides,
    build_pptx_from_text,
    pptx_available,
)

pytestmark = pytest.mark.skipif(
    not pptx_available(), reason="python-pptx not installed"
)


def _slide_titles(path: Path) -> list[str]:
    from pptx import Presentation

    prs = Presentation(str(path))
    return [s.shapes.title.text if s.shapes.title else "" for s in prs.slides]


class TestSplit:
    def test_explicit_slide_markers(self):
        text = (
            "Слайд 1 — Титульный\nПрогноз погоды\nСлайд 2 — 5 июня\n"
            "• Куртка\n• Зонт\nСлайд 3 — 6 июня\nЯсно"
        )
        slides = _split_into_slides(text)
        assert [t for t, _ in slides] == ["Титульный", "5 июня", "6 июня"]
        assert slides[1][1] == ["Куртка", "Зонт"]  # bullet glyphs stripped

    def test_english_slide_markers(self):
        slides = _split_into_slides("Slide 1: Intro\nhi\nSlide 2: Body\nfoo")
        assert [t for t, _ in slides] == ["Intro", "Body"]

    def test_markdown_headers_as_slides(self):
        slides = _split_into_slides("## First\na\nb\n## Second\nc")
        assert [t for t, _ in slides] == ["First", "Second"]
        assert slides[0][1] == ["a", "b"]

    def test_no_markers_falls_back_to_blocks(self):
        slides = _split_into_slides("Title A\nbody a\n\nTitle B\nbody b")
        assert [t for t, _ in slides] == ["Title A", "Title B"]

    def test_unstructured_single_slide(self):
        slides = _split_into_slides("just one line of text")
        assert len(slides) == 1
        assert slides[0][1] == ["just one line of text"]

    def test_empty_text_never_crashes(self):
        slides = _split_into_slides("")
        assert len(slides) == 1


class TestBuild:
    def test_builds_real_pptx(self, tmp_path):
        text = "Слайд 1 — A\n• x\nСлайд 2 — B\n• y"
        dest = tmp_path / "deck.pptx"
        out = build_pptx_from_text(text, dest)
        assert out == dest and dest.exists() and dest.stat().st_size > 0
        assert _slide_titles(dest) == ["A", "B"]

    def test_optional_title_slide_prepended(self, tmp_path):
        dest = tmp_path / "deck.pptx"
        build_pptx_from_text("Слайд 1 — Body\nx", dest, title="My Deck")
        titles = _slide_titles(dest)
        assert titles[0] == "My Deck"  # title slide first
        assert "Body" in titles

    def test_creates_parent_dirs(self, tmp_path):
        dest = tmp_path / "nested" / "sub" / "deck.pptx"
        build_pptx_from_text("Slide 1: X\nhi", dest)
        assert dest.exists()

    def test_overwrites_existing(self, tmp_path):
        dest = tmp_path / "deck.pptx"
        dest.write_text("stale")
        build_pptx_from_text("Slide 1: X\nhi", dest)
        # a real pptx is a zip — first bytes are PK
        assert dest.read_bytes()[:2] == b"PK"


class TestSkillOutputText:
    """ChatScreen._skill_output_text picks the right text to build from."""

    @staticmethod
    def _step(text, *, no_file=False):
        from types import SimpleNamespace

        data = {"no_output_file": True} if no_file else {}
        return SimpleNamespace(result=text, result_data=data)

    def test_prefers_no_output_file_step(self):
        from care.screens.chat import ChatScreen

        result = type("R", (), {"step_results": [
            self._step("forecast data"),
            self._step("Слайд 1 — Title\n...", no_file=True),
        ]})()
        assert "Слайд 1" in ChatScreen._skill_output_text(result)

    def test_falls_back_to_last_nonempty(self):
        from care.screens.chat import ChatScreen

        result = type("R", (), {"step_results": [
            self._step("first"),
            self._step("last"),
        ]})()
        assert ChatScreen._skill_output_text(result) == "last"

    def test_empty_when_no_steps(self):
        from care.screens.chat import ChatScreen

        result = type("R", (), {"step_results": []})()
        assert ChatScreen._skill_output_text(result) == ""
