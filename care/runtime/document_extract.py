"""Plain-text extraction from office / rich-text documents (@-file refs).

CARE's chat surface lets users attach files with ``@<path>``. Plain-text
files embed verbatim, PDFs go through ``pypdf``, images are base64'd for
vision models. This module covers the **binary office documents** in
between — Word (``.docx``), PowerPoint (``.pptx``), Excel
(``.xlsx``/``.xlsm``), OpenDocument (``.odt``/``.odp``/``.ods``) and Rich
Text (``.rtf``) — extracting a plain-text projection that MAGE can read.

Each extractor lazily imports its backing library (the CARE house style: a
missing dependency surfaces as a friendly install hint, not an
``ImportError`` at startup). Dispatch is by file extension via
:func:`extract_document_text`; :data:`DOCUMENT_EXTENSIONS` is the set the
chat surface routes here. Legacy binary formats (``.doc``/``.ppt``/``.xls``)
get a targeted "re-save as …" hint instead of a generic binary-file error.
"""

from __future__ import annotations

from pathlib import Path
from typing import Callable


class DocumentExtractionError(RuntimeError):
    """Extraction failed — the message carries a user-facing reason/hint."""


def _missing(pip_name: str, ext: str) -> "DocumentExtractionError":
    return DocumentExtractionError(
        f"`{pip_name}` isn't installed — needed to read {ext} files. "
        f"Install it with `pip install {pip_name}`."
    )


def _extract_docx(path: Path) -> str:
    """Word: paragraphs first, then table rows as ``a | b | c``."""
    try:
        import docx  # python-docx
    except ImportError as exc:
        raise _missing("python-docx", ".docx") from exc
    document = docx.Document(str(path))
    parts: list[str] = [
        para.text.strip() for para in document.paragraphs if para.text.strip()
    ]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    """PowerPoint: one ``--- slide N ---`` block per slide, text frames +
    tables flattened in shape order."""
    try:
        from pptx import Presentation  # python-pptx
    except ImportError as exc:
        raise _missing("python-pptx", ".pptx") from exc
    presentation = Presentation(str(path))
    slides_out: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        lines.append(" | ".join(cells))
        if lines:
            slides_out.append(f"--- slide {index} ---\n" + "\n".join(lines))
    return "\n\n".join(slides_out)


def _extract_xlsx(path: Path) -> str:
    """Excel: one ``--- sheet <name> ---`` block per worksheet, each
    non-empty row as ``a | b | c`` (trailing empty cells trimmed)."""
    try:
        import openpyxl
    except ImportError as exc:
        raise _missing("openpyxl", ".xlsx") from exc
    # read_only keeps memory bounded on big sheets; data_only reads the
    # last-computed value instead of the formula text.
    workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    try:
        sheets_out: list[str] = []
        for sheet in workbook.worksheets:
            rows_out: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = ["" if value is None else str(value) for value in row]
                while cells and not cells[-1].strip():
                    cells.pop()
                if cells:
                    rows_out.append(" | ".join(cells))
            if rows_out:
                sheets_out.append(
                    f"--- sheet {sheet.title} ---\n" + "\n".join(rows_out)
                )
        return "\n\n".join(sheets_out)
    finally:
        workbook.close()


def _extract_odf(path: Path) -> str:
    """OpenDocument text/presentation/spreadsheet: flatten every
    ``text:p`` paragraph (covers body text, slide text and cell text)."""
    try:
        from odf import teletype, text
        from odf.opendocument import load
    except ImportError as exc:
        raise _missing("odfpy", path.suffix.lower()) from exc
    document = load(str(path))
    lines = [
        teletype.extractText(paragraph)
        for paragraph in document.getElementsByType(text.P)
    ]
    return "\n".join(line for line in lines if line.strip())


def _extract_rtf(path: Path) -> str:
    """Rich Text Format: strip control words to plain text."""
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError as exc:
        raise _missing("striprtf", ".rtf") from exc
    # RTF is 7-bit ASCII with escape sequences; latin-1 never raises and
    # round-trips the high bytes for striprtf's own \'xx decoding.
    raw = path.read_bytes().decode("latin-1", errors="ignore")
    return rtf_to_text(raw)


_EXTRACTORS: dict[str, Callable[[Path], str]] = {
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".xlsm": _extract_xlsx,
    ".odt": _extract_odf,
    ".odp": _extract_odf,
    ".ods": _extract_odf,
    ".rtf": _extract_rtf,
}

#: Extensions the chat surface routes to :func:`extract_document_text`.
DOCUMENT_EXTENSIONS: frozenset[str] = frozenset(_EXTRACTORS)

# Legacy binary formats we can't read without heavyweight native tooling —
# surfaced as a targeted "convert to X" hint rather than a generic
# binary-file warning.
_LEGACY_HINTS: dict[str, str] = {
    ".doc": ".docx",
    ".ppt": ".pptx",
    ".xls": ".xlsx",
}

#: Extensions the chat surface should route to :func:`extract_document_text`:
#: the supported set *plus* the legacy binaries (which yield a "re-save
#: as …" hint instead of falling through to a generic binary-file error).
ROUTABLE_EXTENSIONS: frozenset[str] = DOCUMENT_EXTENSIONS | frozenset(_LEGACY_HINTS)


def extract_document_text(path: Path) -> str:
    """Extract a plain-text projection of *path*, dispatched by extension.

    Raises :class:`DocumentExtractionError` for unsupported extensions, a
    missing backing library, or an unparseable file — the message is
    user-facing (the chat surface shows it verbatim).
    """
    suffix = path.suffix.lower()
    extractor = _EXTRACTORS.get(suffix)
    if extractor is None:
        modern = _LEGACY_HINTS.get(suffix)
        if modern:
            raise DocumentExtractionError(
                f"legacy `{suffix}` files aren't supported — re-save as "
                f"`{modern}` and try again."
            )
        raise DocumentExtractionError(f"unsupported document type `{suffix}`.")
    try:
        return extractor(path)
    except DocumentExtractionError:
        raise
    except Exception as exc:  # noqa: BLE001 — surface any parser error as a hint
        raise DocumentExtractionError(
            f"couldn't parse {suffix} document: {exc}"
        ) from exc
