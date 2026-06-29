"""Deterministic ``.pptx`` builder — the model-independent fallback.

When a file-producing AgentSkill finishes without writing a file (a weak model
couldn't drive the sandbox's ``run_script``/``write_file`` tools, or the
sandbox/python-pptx step failed), the *content* the user asked for has still
been generated — it's sitting in the step's text output as slide descriptions.
This module turns that text into a real ``.pptx`` in-process with ``python-pptx``:
no LLM, no Docker, no network — so the user always gets a file.

The parser is deliberately forgiving: it recognises explicit slide markers
(``Слайд 3 — …`` / ``Slide 3: …`` / Markdown ``##`` headings), bullet glyphs
(``• - * –``), and falls back to blank-line-separated blocks when the text has
no structure at all. It never raises on weird input — worst case it emits a
single content slide with the raw text.

Public API:
    build_pptx_from_text(text, dest, *, title=None) -> Path
    pptx_available() -> bool
"""

from __future__ import annotations

import re
from pathlib import Path

__all__ = ["build_pptx_from_text", "pptx_available"]

# A slide boundary: optional Markdown hashes, then "Слайд"/"Slide"/"Lecture"
# + a number, then an optional dash/colon and the slide's own title.
_SLIDE_HEADER = re.compile(
    r"^\s*#{0,3}\s*(?:слайд|slide|лист|page)\s*\#?\s*\d+\s*[—–\-:.)]*\s*(.*)$",
    re.IGNORECASE,
)
# A plain Markdown heading also starts a new slide.
_MD_HEADER = re.compile(r"^\s*#{1,3}\s+(.*\S)\s*$")
# Leading bullet glyphs / list markers to strip from a body line.
_BULLET_PREFIX = re.compile(r"^\s*(?:[•·▪◦‣*\-–—]|\d+[.)])\s+")

_MAX_SLIDES = 80
_MAX_BULLETS_PER_SLIDE = 18
_MAX_BULLET_CHARS = 240


def pptx_available() -> bool:
    """Whether ``python-pptx`` is importable (the dep ships by default, but a
    minimal install may have dropped it)."""
    try:
        import pptx  # noqa: F401
    except Exception:  # noqa: BLE001
        return False
    return True


def _clean_title(text: str, *, fallback: str) -> str:
    t = (text or "").strip().strip("*#").strip()
    return t or fallback


def _split_into_slides(text: str) -> list[tuple[str, list[str]]]:
    """Parse ``text`` into ``[(title, [bullet, …]), …]``.

    Prefers explicit slide markers; falls back to Markdown headings; finally to
    blank-line-separated blocks (first line = title). Always returns ≥1 slide.
    """
    lines = (text or "").replace("\r\n", "\n").split("\n")

    slides: list[tuple[str, list[str]]] = []
    cur_title: str | None = None
    cur_body: list[str] = []
    saw_marker = False

    def flush() -> None:
        if cur_title is not None or cur_body:
            title = _clean_title(cur_title or "", fallback=(cur_body[0] if cur_body else "Слайд"))
            body = cur_body if cur_title is not None else cur_body[1:]
            slides.append((title, [b for b in body if b][:_MAX_BULLETS_PER_SLIDE]))

    for raw in lines:
        line = raw.rstrip()
        m = _SLIDE_HEADER.match(line) or _MD_HEADER.match(line)
        if m:
            saw_marker = True
            flush()
            cur_title = m.group(1).strip()
            cur_body = []
            continue
        stripped = line.strip()
        if not stripped:
            continue
        bullet = _BULLET_PREFIX.sub("", stripped).strip()
        if bullet:
            cur_body.append(bullet[:_MAX_BULLET_CHARS])
    flush()

    if saw_marker and slides:
        return slides[:_MAX_SLIDES]

    # No markers: fall back to blank-line-separated blocks.
    blocks = [b.strip() for b in re.split(r"\n\s*\n", text or "") if b.strip()]
    if len(blocks) > 1:
        out: list[tuple[str, list[str]]] = []
        for block in blocks[:_MAX_SLIDES]:
            blines = [
                _BULLET_PREFIX.sub("", ln.strip()).strip()
                for ln in block.split("\n") if ln.strip()
            ]
            if not blines:
                continue
            out.append((blines[0][:_MAX_BULLET_CHARS], [b[:_MAX_BULLET_CHARS] for b in blines[1:][:_MAX_BULLETS_PER_SLIDE]]))
        if out:
            return out

    # Truly unstructured: one slide with the whole text as bullets.
    body = [
        _BULLET_PREFIX.sub("", ln.strip()).strip()
        for ln in (text or "").split("\n") if ln.strip()
    ]
    return [("Презентация", body[:_MAX_BULLETS_PER_SLIDE])]


def build_pptx_from_text(text: str, dest: Path, *, title: str | None = None) -> Path:
    """Build a ``.pptx`` at ``dest`` from slide-descriptive ``text``.

    Parameters
    ----------
    text:
        The generated slide content (e.g. the skill step's text output).
    dest:
        Output path. Parent dirs are created; an existing file is overwritten.
    title:
        Optional deck title for the leading title slide. When omitted, the first
        parsed slide carries the deck.

    Returns the written ``dest``. Raises ``RuntimeError`` if ``python-pptx``
    isn't installed; otherwise builds a best-effort deck (never raises on odd
    text).
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except Exception as exc:  # noqa: BLE001
        raise RuntimeError(
            "python-pptx is not installed; cannot build the .pptx fallback"
        ) from exc

    slides = _split_into_slides(text)
    prs = Presentation()

    # Optional leading title slide (layout 0 = Title).
    if title and title.strip():
        s = prs.slides.add_slide(prs.slide_layouts[0])
        s.shapes.title.text = title.strip()[:200]
        if len(s.placeholders) > 1:
            s.placeholders[1].text = "Сгенерировано CARE"

    content_layout = prs.slide_layouts[1]  # Title + Content
    for slide_title, bullets in slides:
        s = prs.slides.add_slide(content_layout)
        s.shapes.title.text = (slide_title or "Слайд")[:200]
        body = s.placeholders[1].text_frame if len(s.placeholders) > 1 else None
        if body is None:
            continue
        if bullets:
            body.text = bullets[0]
            for b in bullets[1:]:
                p = body.add_paragraph()
                p.text = b
                p.level = 0
            for para in body.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(18)

    dest = Path(dest)
    dest.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dest))
    return dest
